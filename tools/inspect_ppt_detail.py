# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt

path = r"presentation/Genomic Encoding via Border Optimization (GERM-BO).pptx"
prs = Presentation(path)
print(f"Slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides, 1):
    print(f"\n=== Slide {i} ===")
    layout = slide.slide_layout
    print(f"  Layout: {layout.name}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f"  [{shape.name}]")
                for para in shape.text_frame.paragraphs:
                    runs_info = []
                    for run in para.runs:
                        font = run.font
                        runs_info.append(
                            f"'{run.text}' size={font.size} bold={font.bold} color={font.color.rgb if font.color and font.color.type else None}"
                        )
                    print(f"    P: {para.text[:120]}")
                    if runs_info:
                        print(f"      runs: {runs_info[:3]}")
        elif shape.shape_type == 13:
            print(f"  [PICTURE] {shape.name} at ({shape.left.inches:.2f}, {shape.top.inches:.2f})")
