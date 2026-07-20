from pptx import Presentation

path = r"presentation/Genomic Encoding via Border Optimization (GERM-BO).pptx"
prs = Presentation(path)
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
for i, slide in enumerate(prs.slides, 1):
    print(f"\n=== Slide {i} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().replace("\n", " | ")
            if text:
                print(f"  [{shape.name}] {text[:300]}")
        elif shape.shape_type == 13:
            print(f"  [PICTURE] {shape.name}")
