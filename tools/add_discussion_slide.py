# -*- coding: utf-8 -*-
"""Build or refresh the Discussion slide in the GERM-BO presentation."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_PATH = Path("presentation/Genomic Encoding via Border Optimization (GERM-BO).pptx")
DISCUSSION_INDEX = 5

C_TITLE = RGBColor(0x0F, 0x17, 0x2A)
C_ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
C_SECTION = RGBColor(0x03, 0x69, 0xA1)
C_BODY = RGBColor(0x47, 0x55, 0x69)
C_EMPH = RGBColor(0x10, 0xB9, 0x81)
C_MUTED = RGBColor(0x64, 0x74, 0x8B)

# Each bullet: (keyword, detail)
ADVANTAGES = [
    ("\u673a\u5236\u660e\u786e", "\u8865\u507f\u8fb9\u754c\u8109\u51b2\u6291\u5236\uff0c\u65e0 LoRA \u989d\u5916\u53c2\u6570"),
    ("\u526a\u63a5\u5b9e\u6d4b", "Macro-F1 +11.8 pp\uff08DNABERT-2\uff0c5 seeds\uff09"),
    ("\u53d7\u63a7\u53ef\u8fc1\u79fb", "border_hard +11 pp\uff1boracle \u4e0b NT/Hyena \u5747\u63d0\u5347"),
]

SCENARIOS = [
    ("\u526a\u63a5\u4f4d\u70b9\u8bc6\u522b", "\u7c7b\u522b\u6781\u5ea6\u4e0d\u5e73\u8861\u3001\u9700\u6062\u590d\u4f9b/\u53d7\u4f53\u4fe1\u53f7"),
    ("border-aware \u4efb\u52a1", "\u6709\u5143\u6570\u636e\u6216\u53ef\u4f30\u8ba1\u8fb9\u754c\u5206\u6570\u7684\u8c03\u63a7\u573a\u666f"),
    ("\u8f7b\u91cf\u6a21\u578b\u9002\u914d", "\u5355\u5361\u5feb\u901f\u5fae\u8c03\uff0c\u6539\u52a8\u6700\u5c0f\u5316"),
]

LIMITATION = (
    "\u5c40\u9650\uff1a\u672a\u8d85\u8d8a 3-mer SVM\uff1b"
    "promoter/chromatin label-free \u589e\u76ca\u63a5\u8fd1 null\uff1b"
    "label-free \u8de8 backbone \u4ecd\u5f85\u9a8c\u8bc1"
)


def _set_run(run, text, *, size_pt, bold=False, color=C_BODY):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"


def _clear_text_shapes(slide):
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            shape._element.getparent().remove(shape._element)


def _add_title_banner(slide):
    box = slide.shapes.add_textbox(Inches(0.83), Inches(0.52), Inches(12.47), Inches(0.50))
    p = box.text_frame.paragraphs[0]
    _set_run(p.add_run(), "Discussion: ", size_pt=30, bold=False, color=C_TITLE)
    _set_run(
        p.add_run(),
        "\u4f18\u52bf\u4e0e\u5e94\u7528",
        size_pt=30,
        bold=False,
        color=C_ACCENT,
    )


def _add_bullets(slide, left, top, width, section_title, items, *, keyword_color=C_EMPH):
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.38)
    )
    _set_run(
        title_box.text_frame.paragraphs[0].add_run(),
        section_title,
        size_pt=22,
        bold=True,
        color=C_SECTION,
    )

    body_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.50), Inches(width), Inches(2.6)
    )
    btf = body_box.text_frame
    btf.word_wrap = True
    for i, (keyword, detail) in enumerate(items):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.level = 0
        p.space_before = Pt(14 if i else 0)
        p.space_after = Pt(6)
        _set_run(p.add_run(), "\u2022 ", size_pt=20, bold=False, color=C_BODY)
        _set_run(p.add_run(), keyword, size_pt=20, bold=True, color=keyword_color)
        _set_run(p.add_run(), "  \u2014  ", size_pt=20, bold=False, color=C_MUTED)
        _set_run(p.add_run(), detail, size_pt=20, bold=False, color=C_BODY)


def _add_limitation_bar(slide):
    box = slide.shapes.add_textbox(Inches(0.83), Inches(5.35), Inches(11.67), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 1  # center
    _set_run(p.add_run(), LIMITATION, size_pt=15, bold=False, color=C_MUTED)


def build_discussion_slide(slide):
    _clear_text_shapes(slide)
    _add_title_banner(slide)
    _add_bullets(
        slide,
        0.62,
        1.30,
        5.85,
        "\u6838\u5fc3\u4f18\u52bf",
        ADVANTAGES,
        keyword_color=C_EMPH,
    )
    _add_bullets(
        slide,
        6.85,
        1.30,
        5.85,
        "\u9002\u7528\u573a\u666f",
        SCENARIOS,
        keyword_color=C_ACCENT,
    )
    _add_limitation_bar(slide)


def insert_discussion_slide(prs, index=DISCUSSION_INDEX):
    blank_layout = next(
        (layout for layout in prs.slide_layouts if layout.name.upper() == "BLANK"),
        prs.slide_layouts[6],
    )
    prs.slides.add_slide(blank_layout)
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(index, new_id)
    build_discussion_slide(prs.slides[index])


def main():
    prs = Presentation(str(PPT_PATH))
    if len(prs.slides) > DISCUSSION_INDEX:
        build_discussion_slide(prs.slides[DISCUSSION_INDEX])
    else:
        insert_discussion_slide(prs)
    prs.save(str(PPT_PATH))
    print(f"Updated slide {DISCUSSION_INDEX + 1} in {PPT_PATH}")


if __name__ == "__main__":
    main()
