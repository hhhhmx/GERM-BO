# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

path = r"presentation/Genomic Encoding via Border Optimization (GERM-BO).pptx"
prs = Presentation(path)
for i, slide in enumerate(prs.slides, 1):
    print(f"\n=== Slide {i} ===")
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f"  {shape.name}: left={shape.left.inches:.2f} top={shape.top.inches:.2f} w={shape.width.inches:.2f} h={shape.height.inches:.2f}")
            print(f"    text={shape.text_frame.text[:80]!r}")
